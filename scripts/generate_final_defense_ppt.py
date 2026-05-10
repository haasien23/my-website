# -*- coding: utf-8 -*-
"""Generate beautiful graduation defense PPT"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn

PROJECT_ROOT = r"D:\biushe"
OUTPUT_DIR = os.path.join(os.path.expanduser("~"), "Desktop", "答辩ppt")
FIGURES_DIR = os.path.join(PROJECT_ROOT, "output", "figures")
LOGO_PATH = os.path.join(PROJECT_ROOT, "output", "ppt_visual_scheme", "jxust-logo-cropped.png")
LOGO_ALT = os.path.join(PROJECT_ROOT, "output", "ppt_visual_scheme", "jxust-logo.jpg")
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "基于集群的多任务协同态势感知平台_答辩PPT.pptx")

os.makedirs(OUTPUT_DIR, exist_ok=True)

# Colors
C_DARK   = RGBColor(0x0D, 0x1B, 0x3E)
C_MID    = RGBColor(0x1B, 0x3A, 0x6B)
C_ACCENT = RGBColor(0xD4, 0xA8, 0x43)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF0, 0xF2, 0xF8)
C_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY   = RGBColor(0x6B, 0x72, 0x84)
C_RED    = RGBColor(0xDC, 0x35, 0x35)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_BLUE_L = RGBColor(0x3B, 0x6F, 0xC6)
C_CARD   = RGBColor(0xE8, 0xEC, 0xF4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def _logo_path():
    if os.path.exists(LOGO_PATH):
        return LOGO_PATH
    return LOGO_ALT

def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_logo(slide, left=None, top=None, width=None):
    if left is None:
        width = width or Inches(1.15)
        left = SLIDE_W - width - Inches(0.4)
    if top is None:
        top = Inches(0.2)
    width = width or Inches(1.15)
    lp = _logo_path()
    if os.path.exists(lp):
        slide.shapes.add_picture(lp, left, top, width=width)

def add_bg_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=None, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
    if color is None:
        color = C_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return txBox, tf

def add_para(tf, text, size=14, color=None, bold=False, alignment=PP_ALIGN.LEFT,
             space_before=0, space_after=4, first=False):
    if color is None:
        color = C_TEXT
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def add_card(slide, left, top, width, height, color=None):
    if color is None:
        color = C_CARD
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.05
    except:
        pass
    return shape

def add_line(slide, x1, y1, x2, y2, color=None, width=1.5):
    if color is None:
        color = C_MID
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector

def slide_number(slide, num):
    add_textbox(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.4), Inches(0.6), Inches(0.3),
                str(num), font_size=10, color=C_GRAY, alignment=PP_ALIGN.RIGHT)

def section_header(slide, title, subtitle=None):
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.05), C_DARK)
    add_textbox(slide, Inches(0.7), Inches(0.15), Inches(10), Inches(0.55),
                title, font_size=28, color=C_WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(0.62), Inches(10), Inches(0.35),
                    subtitle, font_size=13, color=RGBColor(0xBB, 0xC2, 0xD8))
    add_bg_rect(slide, Inches(0.7), Inches(0.95), Inches(1.8), Inches(0.04), C_ACCENT)
    add_logo(slide)

def add_figure(slide, fig_name, left, top, width=None, height=None):
    path = os.path.join(FIGURES_DIR, fig_name)
    if os.path.exists(path):
        if width and height:
            return slide.shapes.add_picture(path, left, top, width=width, height=height)
        elif width:
            return slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            return slide.shapes.add_picture(path, left, top, height=height)
        else:
            return slide.shapes.add_picture(path, left, top, width=Inches(8))
    return None

def build_slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_LIGHT

print("Helper functions defined OK")

# ============ SLIDE 1: COVER ============
def build_slide_01():
    slide = add_blank_slide()
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_DARK)
    # Decorative triangle
    deco = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.5), Inches(0), Inches(5), Inches(5))
    deco.fill.solid(); deco.fill.fore_color.rgb = C_MID; deco.line.fill.background()
    deco.rotation = 180.0
    # Gold accent line
    add_bg_rect(slide, Inches(1.3), Inches(2.2), Inches(0.06), Inches(2.0), C_ACCENT)
    # Title
    add_textbox(slide, Inches(1.7), Inches(2.0), Inches(10), Inches(1.2),
                "基于集群的多任务协同态势感知平台", font_size=38, color=C_WHITE, bold=True)
    add_textbox(slide, Inches(1.7), Inches(3.2), Inches(10), Inches(0.6),
                "Cluster-based Multi-task Collaborative Situation Awareness Platform",
                font_size=17, color=RGBColor(0xAA, 0xB4, 0xCC))
    # Tags
    tags = ["RADS动态子群选择", "复杂环境仿真建模", "Web三维可视化", "多策略对比实验"]
    for i, tag in enumerate(tags):
        add_textbox(slide, Inches(1.7) + Inches(i * 2.5), Inches(3.8), Inches(2.2), Inches(0.35),
                    tag, font_size=11, color=C_ACCENT, bold=True)
    # Info
    add_textbox(slide, Inches(1.7), Inches(4.8), Inches(8), Inches(0.4),
                "答辩人: 贺小双    专业: 计算机科学与技术    班级: 计算机224班",
                font_size=14, color=RGBColor(0xCC, 0xD2, 0xE0))
    add_textbox(slide, Inches(1.7), Inches(5.2), Inches(8), Inches(0.4),
                "指导教师: XXX 教授    江西理工大学 信息工程学院",
                font_size=14, color=RGBColor(0xCC, 0xD2, 0xE0))
    add_logo(slide, width=Inches(1.4))
    add_bg_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), RGBColor(0x08, 0x12, 0x2E))
    add_textbox(slide, Inches(1.3), Inches(7.0), Inches(6), Inches(0.35),
                "江西理工大学 · 信息工程学院 · 本科毕业设计答辩",
                font_size=10, color=C_GRAY)
    return slide

# ============ SLIDE 2: TOC ============
def build_slide_02():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "汇报提纲", "Presentation Outline")
    items = [
        ("01", "研究背景与问题", "无人机集群应用场景、复杂环境约束与现有方法不足"),
        ("02", "研究目标与技术路线", "场景建模 -> 算法设计 -> 平台实现 -> 实验验证"),
        ("03", "系统设计与架构", "前后端分离架构、三大层次、核心模块划分"),
        ("04", "RADS 动态子群选择算法", "目标需求评估、节点效用计算、动态预算分配"),
        ("05", "平台实现与展示", "三维态势沙盘、多策略对比、任务级分析、结果导出"),
        ("06", "实验设计与结果分析", "四组实验、多指标评价、综合对比"),
        ("07", "总结与展望", "创新点、结论、不足与改进方向"),
    ]
    start_y = Inches(1.4)
    for i, (num, title, desc) in enumerate(items):
        y = start_y + Inches(i * 0.78)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.5), Inches(0.5))
        circle.fill.solid(); circle.fill.fore_color.rgb = C_DARK
        circle.line.fill.background()
        tf = circle.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(14); p.font.color.rgb = C_WHITE
        p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        add_textbox(slide, Inches(1.8), y, Inches(5), Inches(0.35),
                    title, font_size=17, color=C_DARK, bold=True)
        add_textbox(slide, Inches(1.8), y + Inches(0.35), Inches(9), Inches(0.3),
                    desc, font_size=11, color=C_GRAY)
        if i < len(items) - 1:
            add_line(slide, Inches(1.8), y + Inches(0.68), Inches(11.5), y + Inches(0.68),
                     color=RGBColor(0xD0, 0xD5, 0xE0), width=0.8)
    slide_number(slide, 2)
    return slide

print("build_slide_01 and build_slide_02 defined OK")

# ============ SLIDE 13: DEFAULT SCENE RESULTS ============
def build_slide_13():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "默认场景实验结果", "Default Scene Results -- 24机 x 6目标 x 34障碍 x 晴空")
    metrics = [
        ("严格成功率", "RADS: 70.8%", "随机: 16.7%", "全量: 72.9%",
         "RADS接近全量,\n远超随机"),
        ("平均定位误差", "RADS: 9.52m", "随机: 47.61m", "全量: 8.07m",
         "RADS误差与全量\n在同一水平"),
        ("累计能耗", "RADS: 1120.4", "随机: 1109.8", "全量: 2138.5",
         "RADS比全量\n降低 47.6%"),
        ("平均派机数量", "RADS: 9.4架", "随机: 9.4架", "全量: 20.5架",
         "RADS比全量\n降低 54.1%"),
    ]
    for i, (metric, rads, rand, full, note) in enumerate(metrics):
        x = Inches(0.3) + Inches(i * 3.25)
        add_card(slide, x, Inches(1.3), Inches(3.05), Inches(3.3), C_WHITE)
        add_textbox(slide, x + Inches(0.2), Inches(1.4), Inches(2.65), Inches(0.3),
                    metric, font_size=14, color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        colors = [C_GREEN, C_RED, C_BLUE_L]
        labels = ["RADS", "随机", "全量"]
        values = [rads, rand, full]
        for j in range(3):
            mx = x + Inches(0.1) + Inches(j * 1.0)
            add_textbox(slide, mx, Inches(1.8), Inches(0.9), Inches(0.25),
                        labels[j], font_size=9, color=C_GRAY, alignment=PP_ALIGN.CENTER)
            add_textbox(slide, mx, Inches(2.05), Inches(0.9), Inches(0.5),
                        values[j], font_size=12, color=colors[j], bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(2.9), Inches(2.65), Inches(0.4),
                    note, font_size=9, color=C_GRAY, alignment=PP_ALIGN.CENTER)
        add_bg_rect(slide, x + Inches(0.2), Inches(2.8), Inches(2.65), Inches(0.02), C_ACCENT)
    # Key findings
    add_card(slide, Inches(0.5), Inches(4.85), Inches(12.1), Inches(2.3), C_DARK)
    add_textbox(slide, Inches(0.9), Inches(5.0), Inches(11), Inches(0.35),
                "关键发现", font_size=16, color=C_ACCENT, bold=True)
    findings = [
        "RADS 严格成功率 70.8%, 接近全量派遣的 72.9%, 同时远超随机派遣的 16.7% -- 说明动态子群选择是有效的",
        "RADS 累计能耗仅约为全量派遣的一半 (47.6%降幅), 平均派机数量仅约 46% -- 以更低资源代价达到接近最优的效果",
        "随机派遣虽然资源投入与RADS相当, 但成功率极低且误差很大 -- 说明派谁去和派多少同样重要",
        "RADS 的成功不仅因为控制了派机规模, 更因为每一步都选择了最适合当前目标需求的无人机成员",
    ]
    for i, f in enumerate(findings):
        add_textbox(slide, Inches(0.9), Inches(5.45) + Inches(i * 0.38), Inches(11.3), Inches(0.35),
                    "> " + f, font_size=11, color=C_WHITE)
    add_figure(slide, "ch8/fig8-7_sr_error_compare.png", Inches(7.5), Inches(1.3), width=Inches(5.0))
    add_figure(slide, "ch8/fig8-8_energy_dispatch_compare.png", Inches(7.5), Inches(3.0), width=Inches(5.0))
    slide_number(slide, 13)
    return slide

# ============ SLIDE 14: TASK SCALE EXPERIMENT ============
def build_slide_14():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "实验一: 任务规模影响", "Task Scale Experiment -- 目标数 6->10")
    add_figure(slide, "ch8/line_task_scale.png", Inches(0.3), Inches(1.2), width=Inches(7.0))
    add_card(slide, Inches(7.8), Inches(1.3), Inches(5.0), Inches(2.2), C_WHITE)
    add_textbox(slide, Inches(8.1), Inches(1.45), Inches(4.5), Inches(0.35),
                "严格成功率变化", font_size=16, color=C_DARK, bold=True)
    scale_data = [
        "目标数 6->10, 任务负载明显上升",
        "RADS: 70.8% -> 68.0% (仅下降 2.8pp)",
        "随机: 16.7% -> 10.8% (下降 5.9pp)",
        "全量: 72.9% -> 58.3% (下降 14.6pp!)",
    ]
    scale_colors = [C_TEXT, C_GREEN, C_RED, C_BLUE_L]
    for i, (text, color) in enumerate(zip(scale_data, scale_colors)):
        add_textbox(slide, Inches(8.1), Inches(1.9) + Inches(i * 0.38), Inches(4.5), Inches(0.3),
                    "> " + text if i > 0 else text, font_size=12, color=color)
    add_card(slide, Inches(7.8), Inches(3.7), Inches(5.0), Inches(3.5), C_DARK)
    add_textbox(slide, Inches(8.1), Inches(3.85), Inches(4.5), Inches(0.35),
                "实验结论", font_size=16, color=C_ACCENT, bold=True)
    conclusions = [
        "RADS 对任务负载增加表现出良好的适应性, 成功率仅小幅下降",
        "全量派遣在目标增多后资源竞争加剧, 全派策略效率大幅下滑",
        "RADS 通过动态预算机制, 在任务压力增大时自动扩大子群规模",
        "说明 RADS 的资源弹性分配策略面对任务规模变化时更加稳定",
    ]
    for i, c in enumerate(conclusions):
        add_textbox(slide, Inches(8.1), Inches(4.3) + Inches(i * 0.45), Inches(4.5), Inches(0.4),
                    "> " + c, font_size=11, color=C_WHITE)
    slide_number(slide, 14)
    return slide

# ============ SLIDE 15: WEATHER EXPERIMENT ============
def build_slide_15():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "实验二: 天气扰动影响", "Weather Disturbance -- 晴空->薄雾->降雨->雷暴")
    add_figure(slide, "ch8/line_weather.png", Inches(0.3), Inches(1.2), width=Inches(7.8))
    add_card(slide, Inches(8.5), Inches(1.3), Inches(4.3), Inches(2.3), C_WHITE)
    add_textbox(slide, Inches(8.8), Inches(1.45), Inches(3.8), Inches(0.35),
                "严格成功率对比", font_size=16, color=C_DARK, bold=True)
    weather_data = [
        ("晴空", "RADS 70.8%", "全量 72.9%", "随机 16.7%"),
        ("薄雾", "RADS 65.2%", "全量 68.4%", "随机 14.1%"),
        ("降雨", "RADS 48.3%", "全量 55.8%", "随机 10.0%"),
        ("雷暴", "RADS 34.2%", "全量 50.0%", "随机  6.7%"),
    ]
    for i, (weather, rads, full, rand) in enumerate(weather_data):
        y = Inches(1.9) + Inches(i * 0.38)
        add_textbox(slide, Inches(8.8), y, Inches(1.0), Inches(0.3),
                    weather, font_size=11, color=C_DARK, bold=True)
        add_textbox(slide, Inches(9.8), y, Inches(1.4), Inches(0.3),
                    rads, font_size=10, color=C_GREEN)
        add_textbox(slide, Inches(11.0), y, Inches(1.3), Inches(0.3),
                    full, font_size=10, color=C_BLUE_L)
    add_card(slide, Inches(8.5), Inches(3.85), Inches(4.3), Inches(3.4), C_DARK)
    add_textbox(slide, Inches(8.8), Inches(4.0), Inches(3.8), Inches(0.35),
                "实验结论", font_size=16, color=C_ACCENT, bold=True)
    wc = [
        "天气恶化对所有策略都有负面影响, 但 RADS 下降幅度更可控",
        "雷暴条件下, RADS成功率(34.2%)仍显著优于随机(6.7%)",
        "全量派遣在极端天气下保持最高成功率, 但能耗代价巨大",
        "RADS的优势在于性价比: 用更低资源换取可接受的感知效果",
    ]
    for i, c in enumerate(wc):
        add_textbox(slide, Inches(8.8), Inches(4.45) + Inches(i * 0.45), Inches(3.8), Inches(0.4),
                    "> " + c, font_size=11, color=C_WHITE)
    slide_number(slide, 15)
    return slide

# ============ SLIDE 16: OBSTACLE + COMM EXPERIMENTS ============
def build_slide_16():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "实验三&四: 障碍环境与通信条件", "Obstacle & Communication Experiments")
    # Obstacle (left)
    add_card(slide, Inches(0.3), Inches(1.3), Inches(6.3), Inches(5.8), C_WHITE)
    add_textbox(slide, Inches(0.6), Inches(1.45), Inches(5.8), Inches(0.35),
                "实验三: 障碍复杂度 (34->50个)", font_size=16, color=C_DARK, bold=True)
    add_figure(slide, "ch8/line_obstacle.png", Inches(0.5), Inches(1.9), width=Inches(5.8))
    obs = [
        "障碍增加使绕障路径变长, 可见性降低",
        "RADS 成功率: 70.8% -> 68.8% (仅降2.0pp)",
        "路径代价和视线判断提升调度合理性",
        "RADS自动排除被遮挡的无人机, 减少无效派机",
    ]
    for i, o in enumerate(obs):
        add_textbox(slide, Inches(0.6), Inches(4.7) + Inches(i * 0.35), Inches(5.8), Inches(0.3),
                    "> " + o, font_size=11, color=C_TEXT)
    # Communication (right)
    add_card(slide, Inches(6.9), Inches(1.3), Inches(6.1), Inches(5.8), C_WHITE)
    add_textbox(slide, Inches(7.2), Inches(1.45), Inches(5.5), Inches(0.35),
                "实验四: 通信条件变化", font_size=16, color=C_DARK, bold=True)
    comm_items = [
        ("通信半径 300->400m", "RADS: 70.8% -> 81.7% (+10.9pp)", C_GREEN),
        ("丢包率 0% -> 10%",    "RADS: 77.9% -> 65.0% (-12.9pp)", C_RED),
        ("延迟 0 -> 2步",       "RADS: 70.8% -> 60.0% (-10.8pp)", RGBColor(0xE6, 0x7E, 0x22)),
    ]
    for i, (label, result, color) in enumerate(comm_items):
        y = Inches(1.95) + Inches(i * 1.3)
        bg_color = RGBColor(0xF5, 0xF7, 0xFC) if i % 2 == 0 else C_WHITE
        add_card(slide, Inches(7.2), y, Inches(5.5), Inches(1.1), bg_color)
        add_textbox(slide, Inches(7.4), y + Inches(0.1), Inches(5.1), Inches(0.3),
                    label, font_size=13, color=C_DARK, bold=True)
        add_textbox(slide, Inches(7.4), y + Inches(0.45), Inches(5.1), Inches(0.3),
                    result, font_size=12, color=color)
        add_textbox(slide, Inches(7.4), y + Inches(0.75), Inches(5.1), Inches(0.25),
                    "链路质量直接影响融合权重, 有效观测能否及时送达是协同感知关键",
                    font_size=9, color=C_GRAY)
    comm_notes = [
        "通信半径扩大提升观测到达率, RADS成功率显著上升",
        "丢包率和延迟对融合质量有直接负面影响",
        "RADS通过链路质量评分机制优选通信条件好的节点",
    ]
    for i, n in enumerate(comm_notes):
        add_textbox(slide, Inches(7.2), Inches(5.0) + Inches(i * 0.35), Inches(5.5), Inches(0.3),
                    "> " + n, font_size=11, color=C_TEXT)
    slide_number(slide, 16)
    return slide

# ============ SLIDE 17: COMPREHENSIVE DATA ============
def build_slide_17():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "关键实验数据汇总", "Key Experimental Data Summary")
    add_figure(slide, "ch8/tab8-5_overall_results.png", Inches(0.5), Inches(1.3), width=Inches(7.8))
    # Summary table card
    add_card(slide, Inches(8.8), Inches(1.3), Inches(4.0), Inches(3.5), C_WHITE)
    add_textbox(slide, Inches(9.1), Inches(1.45), Inches(3.5), Inches(0.35),
                "多组实验平均值", font_size=16, color=C_DARK, bold=True)
    avg_data = [
        ("指标", "RADS", "随机", "全量"),
        ("严格成功率", "64.9%", "17.1%", "70.3%"),
        ("平均误差(m)", "10.91", "14.83", "8.66"),
        ("累计能耗", "1060.4", "1056.1", "2143.7"),
        ("平均派机", "8.8", "8.8", "20.4"),
        ("定位达标率", "73.6%", "23.1%", "78.1%"),
        ("信息时效(步)", "0.82", "1.23", "0.68"),
    ]
    for i, row in enumerate(avg_data):
        y = Inches(1.9) + Inches(i * 0.35)
        for j, cell in enumerate(row):
            x = Inches(9.0) + Inches(j * 1.1)
            if i == 0:
                color = C_GRAY
            elif j == 1:
                color = C_GREEN
            elif j == 3:
                color = C_BLUE_L
            elif j == 2:
                color = C_RED
            else:
                color = C_TEXT
            add_textbox(slide, x, y, Inches(1.0), Inches(0.3),
                        cell, font_size=10, color=color, bold=(i == 0 or j == 0))
    # Key takeaway
    add_card(slide, Inches(8.8), Inches(5.1), Inches(4.0), Inches(2.1), C_DARK)
    add_textbox(slide, Inches(9.1), Inches(5.25), Inches(3.5), Inches(0.35),
                "核心结论", font_size=16, color=C_ACCENT, bold=True)
    key_conclusions = [
        "RADS ~ 全量的感知质量",
        "RADS << 全量的资源消耗",
        "RADS >> 随机的稳定性",
        "效果-资源综合最优!",
    ]
    for i, kc in enumerate(key_conclusions):
        bold = (i == 3)
        add_textbox(slide, Inches(9.1), Inches(5.7) + Inches(i * 0.38), Inches(3.5), Inches(0.3),
                    "* " + kc, font_size=14, color=C_WHITE, bold=bold)
    # Bottom charts
    add_figure(slide, "ch8/line_comm_range.png", Inches(0.5), Inches(5.3), width=Inches(2.5))
    add_figure(slide, "ch8/line_packet_loss.png", Inches(3.2), Inches(5.3), width=Inches(2.5))
    add_figure(slide, "ch8/line_comm_delay.png", Inches(5.9), Inches(5.3), width=Inches(2.5))
    slide_number(slide, 17)
    return slide

# ============ SLIDE 18: COMPREHENSIVE ANALYSIS ============
def build_slide_18():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "综合结果分析", "Comprehensive Analysis -- 为什么RADS更好?")
    dims = [
        ("效果维度", Inches(0.5), [
            "严格成功率接近全量派遣 (64.9% vs 70.3%)",
            "定位误差显著优于随机派遣 (10.91m vs 14.83m)",
            "定位达标率与全量差距不大",
            "在4组实验中保持相对稳定的感知质量",
        ]),
        ("效率维度", Inches(4.5), [
            "累计能耗仅为全量派遣的约 49.5%",
            "平均派机数量仅为全量派遣的约 43%",
            "每架无人机承担更具针对性的任务",
            "避免了全员出动带来的大量无效飞行",
        ]),
        ("鲁棒性维度", Inches(8.5), [
            "天气恶化时下降幅度可控 (优于随机4-5倍)",
            "障碍增多时成功率仅降2pp (路径代价自适应)",
            "通信条件变化时有链路质量评分缓冲",
            "任务规模增大时自动扩大子群适应负载",
        ]),
    ]
    for title, x, items in dims:
        add_card(slide, x, Inches(1.3), Inches(3.8), Inches(2.8), C_WHITE)
        add_bg_rect(slide, x, Inches(1.3), Inches(3.8), Inches(0.06), C_ACCENT)
        add_textbox(slide, x + Inches(0.3), Inches(1.5), Inches(3.2), Inches(0.35),
                    title, font_size=17, color=C_DARK, bold=True)
        for i, item in enumerate(items):
            add_textbox(slide, x + Inches(0.3), Inches(1.95) + Inches(i * 0.42), Inches(3.2), Inches(0.35),
                        "/ " + item, font_size=11, color=C_TEXT)
    # Why RADS works
    add_bg_rect(slide, Inches(0.5), Inches(4.4), Inches(12.1), Inches(2.8), C_DARK)
    add_textbox(slide, Inches(0.9), Inches(4.55), Inches(11), Inches(0.35),
                "RADS 有效性的深层原因分析", font_size=17, color=C_ACCENT, bold=True)
    why = [
        "需求感知: 不是每个目标分一样多的无人机, 而是根据优先级、不确定性和运动状态动态评估需求",
        "代价敏感: 不只算直线距离, 而是通过栅格距离图计算绕障路径代价, 避免看起来近但实际绕远的误判",
        "质量筛选: 通过链路质量评分和视线判断, 优先选择通信可靠、视野无遮挡的节点, 减少无效观测",
        "弹性适应: 动态预算机制让派机规模随任务压力自动调整 -- 压力大时多派, 压力小时少派, 避免资源浪费",
        "稳健融合: 多源观测加权融合中抑制异常值和低质量链路观测, 确保即使个别节点失效, 整体结果仍可信",
    ]
    for i, w in enumerate(why):
        add_textbox(slide, Inches(0.9), Inches(5.0) + Inches(i * 0.4), Inches(11.3), Inches(0.35),
                    "> " + w, font_size=11.5, color=C_WHITE)
    slide_number(slide, 18)
    return slide

# ============ SLIDE 19: INNOVATION POINTS ============
def build_slide_19():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "创新点与贡献", "Innovations & Contributions")
    innovations = [
        ("创新点一", "面向复杂约束的统一仿真平台",
         "将无人机、动态目标、障碍物、通信链路、天气扰动和节点故障六类要素纳入统一仿真模型, "
         "使实验环境不再停留在理想化平面场景。平台集参数配置、实验运行、过程回放、任务级分析"
         "和结果导出于一体, 为协同感知研究提供了系统化验证工具。"),
        ("创新点二", "RADS 动态子群选择算法",
         "综合考虑目标优先级、不确定性、运动状态、路径代价、视线关系、链路质量和节点剩余能量"
         "等多维信息, 在每一步自适应地确定无人机子群的规模与成员。与固定派机或简单随机策略相比, "
         "RADS 在复杂环境下实现了更精准的感知资源分配。"),
        ("创新点三", "三维可视化 + 多策略对比 + 任务级分析",
         "平台提供三策略同步对比、三维态势沙盘、回放控制、目标追踪联动和结构化CSV导出功能。"
         "用户不仅能看到结果好不好, 还能从任务级理解为什么好或为什么差, "
         "实现了指标统计 + 场景展示 + 目标级分析的完整验证闭环。"),
        ("创新点四", "多维度综合评价体系",
         "从严格成功率、定位达标率、平均误差、累计能耗、平均派机规模和信息时效六个维度, "
         "对三种策略进行统一评价, 并将实验结果整理为可直接用于论文图表的数据格式, "
         "使实验结论既有数值支撑, 也有可视化佐证。"),
    ]
    for i, (label, title, desc) in enumerate(innovations):
        y = Inches(1.3) + Inches(i * 1.48)
        add_card(slide, Inches(0.5), y, Inches(12.1), Inches(1.3), C_WHITE)
        add_bg_rect(slide, Inches(0.5), y, Inches(0.07), Inches(1.3), C_ACCENT)
        nums = ["[1]", "[2]", "[3]", "[4]"]
        add_textbox(slide, Inches(1.0), y + Inches(0.15), Inches(1.0), Inches(0.4),
                    nums[i], font_size=22, color=C_DARK, bold=True)
        add_textbox(slide, Inches(2.0), y + Inches(0.08), Inches(1.0), Inches(0.25),
                    label, font_size=10, color=C_ACCENT, bold=True)
        add_textbox(slide, Inches(2.0), y + Inches(0.28), Inches(9.8), Inches(0.3),
                    title, font_size=16, color=C_DARK, bold=True)
        add_textbox(slide, Inches(2.0), y + Inches(0.55), Inches(9.8), Inches(0.6),
                    desc, font_size=11, color=C_TEXT)
    slide_number(slide, 19)
    return slide

# ============ SLIDE 20: CONCLUSION & THANKS ============
def build_slide_20():
    slide = add_blank_slide()
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_DARK)
    deco = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.5), Inches(0), Inches(4.2), Inches(4.2))
    deco.fill.solid(); deco.fill.fore_color.rgb = C_MID; deco.line.fill.background()
    deco.rotation = 180.0
    add_logo(slide, width=Inches(1.3))
    add_textbox(slide, Inches(1.3), Inches(1.5), Inches(10), Inches(0.7),
                "总结与展望", font_size=36, color=C_WHITE, bold=True)
    add_bg_rect(slide, Inches(1.3), Inches(2.2), Inches(1.5), Inches(0.04), C_ACCENT)
    add_textbox(slide, Inches(1.3), Inches(2.6), Inches(5), Inches(0.35),
                "已完成的工作", font_size=18, color=C_ACCENT, bold=True)
    conclusions = [
        "实现了面向复杂环境的无人机集群多任务协同态势感知平台",
        "提出了 RADS 动态子群选择算法, 在感知质量与资源消耗之间取得平衡",
        "完成了三策略四组对比实验, 验证了 RADS 的有效性与适应性",
    ]
    for i, c in enumerate(conclusions):
        add_textbox(slide, Inches(1.3), Inches(3.1) + Inches(i * 0.42), Inches(7.5), Inches(0.38),
                    "/ " + c, font_size=13, color=C_WHITE)
    add_textbox(slide, Inches(1.3), Inches(4.6), Inches(5), Inches(0.35),
                "不足与展望", font_size=18, color=C_ACCENT, bold=True)
    futures = [
        "目前仍是任务级仿真, 后续可向更高保真度 (如对接ROS/PX4) 发展",
        "RADS 参数权重目前依赖经验设定, 可引入自适应调参或学习方法",
        "实验梯度可以进一步细化, 覆盖更多中间场景配置",
        "平台可扩展: 接入新的调度策略、天气模型或评价指标",
    ]
    for i, f in enumerate(futures):
        add_textbox(slide, Inches(1.3), Inches(5.1) + Inches(i * 0.38), Inches(7.5), Inches(0.35),
                    "-> " + f, font_size=12, color=RGBColor(0xBB, 0xC2, 0xD8))
    add_textbox(slide, Inches(1.3), Inches(6.5), Inches(10), Inches(0.6),
                "感谢各位老师评阅与指导, 请批评指正!",
                font_size=24, color=C_ACCENT, bold=True)
    add_bg_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), RGBColor(0x08, 0x12, 0x2E))
    add_textbox(slide, Inches(1.3), Inches(7.0), Inches(6), Inches(0.35),
                "江西理工大学 · 信息工程学院 · 计算机224班 · 贺小双",
                font_size=10, color=C_GRAY)
    slide_number(slide, 20)
    return slide

def build_slide_03():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "研究背景", "Research Background -- 为何关注无人机集群协同感知?")
    add_card(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.6), C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(1.45), Inches(5), Inches(0.35),
                "应用场景不断扩展", font_size=18, color=C_DARK, bold=True)
    for i, s in enumerate([
        "安防巡检: 城市、边境、重点设施的大范围巡逻监测",
        "灾害救援: 灾后快速侦察、被困人员搜索、灾情评估",
        "目标侦察: 多目标动态追踪与协同定位",
        "环境监测: 污染扩散、森林防火、水域巡查",
    ]):
        add_textbox(slide, Inches(0.9), Inches(1.95) + Inches(i * 0.45), Inches(5.2), Inches(0.35),
                    "> " + s, font_size=13, color=C_TEXT)
    add_card(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.6), C_WHITE)
    add_textbox(slide, Inches(7.1), Inches(1.45), Inches(5), Inches(0.35),
                "复杂环境的现实挑战", font_size=18, color=C_DARK, bold=True)
    challenges = [
        ("目标动态变化", "目标持续移动、优先级和不确定性实时变化"),
        ("障碍物遮挡", "建筑、地形遮挡影响感知视线与通信链路"),
        ("通信条件受限", "链路丢包、传输延迟、通信半径限制"),
        ("天气与节点故障", "降雨/雷暴扰动、传感器噪声、节点能量衰减与故障"),
    ]
    for i, (c_title, c_desc) in enumerate(challenges):
        y = Inches(1.95) + Inches(i * 0.45)
        add_textbox(slide, Inches(7.2), y, Inches(2.0), Inches(0.3),
                    "! " + c_title, font_size=13, color=C_RED, bold=True)
        add_textbox(slide, Inches(9.2), y, Inches(3.2), Inches(0.3),
                    c_desc, font_size=12, color=C_GRAY)
    add_card(slide, Inches(0.5), Inches(4.15), Inches(12.1), Inches(2.9), C_DARK)
    add_textbox(slide, Inches(0.9), Inches(4.35), Inches(11), Inches(0.4),
                "核心矛盾: 资源有限 vs 需求多变", font_size=20, color=C_ACCENT, bold=True)
    summary = (
        "单架无人机能力有限, 集群化可提升覆盖、并行处理和系统冗余能力。"
        "然而在实际任务中, 目标持续移动、障碍物形成遮挡、通信链路受距离和天气影响、节点存在能量衰减和随机故障, "
        "这些约束使得简单派出所有无人机的策略不再适用。\n"
        "态势感知的关键在于: 对环境状态、目标分布和任务执行情况形成持续、准确、可更新的理解, "
        "在感知质量与资源消耗之间取得合理平衡。"
    )
    add_textbox(slide, Inches(0.9), Inches(4.9), Inches(11.3), Inches(1.8),
                summary, font_size=13, color=C_WHITE)
    slide_number(slide, 3)
    return slide

# ============ SLIDE 4: RESEARCH PROBLEM ============
def build_slide_04():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "研究问题与切入点", "Research Questions -- 本课题要解决什么问题?")
    problems = [
        ("问题一", "多目标并行争夺有限资源",
         "多个目标同时需要感知, 无人机数量有限。\n如何判断每个目标该分配多少资源?\n目标优先级不同、不确定性不同、运动状态不同,\n资源分配不能一刀切。",
         C_RED),
        ("问题二", "复杂环境导致感知条件恶劣",
         "建筑物遮挡视线、通信链路丢包延迟、\n天气降低能见度、节点可能故障。\n理想环境下有效的策略,\n在复杂约束下可能完全失效。",
         RGBColor(0xE6, 0x7E, 0x22)),
        ("问题三", "评价维度单一, 过程不透明",
         "只看最终成功率不够,\n需要同时关注误差、能耗、派机规模、信息时效。\n需要任务级分析能力,\n解释成功或失败的具体原因。",
         RGBColor(0x8E, 0x44, 0xAD)),
    ]
    for i, (label, title, desc, accent) in enumerate(problems):
        x = Inches(0.5) + Inches(i * 4.1)
        add_card(slide, x, Inches(1.4), Inches(3.8), Inches(4.8), C_WHITE)
        add_bg_rect(slide, x, Inches(1.4), Inches(3.8), Inches(0.06), accent)
        add_textbox(slide, x + Inches(0.3), Inches(1.65), Inches(3.2), Inches(0.35),
                    label, font_size=12, color=accent, bold=True)
        add_textbox(slide, x + Inches(0.3), Inches(2.0), Inches(3.2), Inches(0.5),
                    title, font_size=18, color=C_DARK, bold=True)
        add_textbox(slide, x + Inches(0.3), Inches(2.6), Inches(3.2), Inches(3.0),
                    desc, font_size=12, color=C_TEXT)
    add_bg_rect(slide, Inches(0.5), Inches(6.5), Inches(12.1), Inches(0.65), C_DARK)
    add_textbox(slide, Inches(0.9), Inches(6.58), Inches(11.3), Inches(0.5),
                "本课题切入点: 构建统一仿真平台 -> 设计RADS动态子群算法 -> 多策略对比 -> 三维可视化 + 任务级分析",
                font_size=14, color=C_ACCENT, bold=True)
    slide_number(slide, 4)
    return slide

# ============ SLIDE 5: EXISTING METHODS LIMITATIONS ============
def build_slide_05():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "现有方法不足与本课题定位", "Limitations of Existing Approaches & Our Position")
    cols = [
        ("现有方法的问题", [
            "环境建模偏理想化: 忽略障碍、天气、\n链路丢包和节点故障等真实约束",
            "资源分配不够灵活: 固定派机或简单均\n衡, 难以适应目标优先级动态变化",
            "实验验证不完整: 只给数值结果或局部\n示意图, 过程不够透明",
            "评价维度偏窄: 只看单一指标, 不做\n多维度综合权衡",
        ], C_RED),
        ("本课题的改进", [
            "统一复杂场景建模: 障碍物、动态目标、\n天气扰动、通信约束、节点故障一体化",
            "RADS动态子群选择: 每步重新评估目标\n需求和节点效用, 按需调配资源",
            "可视化实验平台: 三维态势沙盘 + 回放\n控制 + 三策略同步对比 + 任务级分析",
            "多维度综合评价: 成功率、误差、能耗、\n派机规模、信息时效统一比较",
        ], C_GREEN),
    ]
    for col_idx, (title, items, accent) in enumerate(cols):
        x = Inches(0.5) + Inches(col_idx * 6.2)
        add_card(slide, x, Inches(1.35), Inches(5.8), Inches(5.8), C_WHITE)
        add_bg_rect(slide, x, Inches(1.35), Inches(5.8), Inches(0.9), accent)
        add_textbox(slide, x + Inches(0.3), Inches(1.5), Inches(5.2), Inches(0.5),
                    title, font_size=18, color=C_WHITE, bold=True)
        for i, item in enumerate(items):
            y = Inches(2.5) + Inches(i * 1.15)
            add_card(slide, x + Inches(0.25), y, Inches(5.3), Inches(1.0), RGBColor(0xF8, 0xF9, 0xFC))
            add_textbox(slide, x + Inches(0.45), y + Inches(0.1), Inches(4.9), Inches(0.8),
                        "> " + item, font_size=12, color=C_TEXT)
    slide_number(slide, 5)
    return slide

# ============ SLIDE 6: RESEARCH OBJECTIVES ============
def build_slide_06():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "研究目标与技术路线", "Research Objectives & Technical Roadmap")
    objectives = [
        ("[1]", "复杂场景建模", "统一描述无人机、目标、\n障碍、天气、通信、故障"),
        ("[2]", "RADS算法设计", "动态子群选择, 兼顾\n目标需求与节点能力"),
        ("[3]", "Web平台实现", "参数配置、三维沙盘、\n回放、对比、导出"),
        ("[4]", "多策略实验验证", "三策略对比、四组实验、\n多维度综合评价"),
    ]
    for i, (icon, title, desc) in enumerate(objectives):
        x = Inches(0.5) + Inches(i * 3.1)
        add_card(slide, x, Inches(1.4), Inches(2.8), Inches(2.0), C_WHITE)
        add_textbox(slide, x + Inches(0.3), Inches(1.55), Inches(2.2), Inches(0.4),
                    icon, font_size=24, color=C_DARK)
        add_textbox(slide, x + Inches(0.3), Inches(1.95), Inches(2.2), Inches(0.35),
                    title, font_size=15, color=C_DARK, bold=True)
        add_textbox(slide, x + Inches(0.3), Inches(2.35), Inches(2.2), Inches(0.8),
                    desc, font_size=11, color=C_GRAY)
    add_bg_rect(slide, Inches(0.5), Inches(3.7), Inches(12.1), Inches(3.5), C_DARK)
    add_textbox(slide, Inches(0.9), Inches(3.85), Inches(11), Inches(0.4),
                "技术路线", font_size=16, color=C_ACCENT, bold=True)
    steps = [
        "需求分析与\n问题定义", "复杂环境\n仿真建模", "RADS算法\n设计与实现",
        "Web前后端\n平台搭建", "三策略\n对比实验", "结果分析\n论文撰写",
    ]
    for i, step in enumerate(steps):
        x = Inches(1.0) + Inches(i * 2.0)
        add_card(slide, x, Inches(4.5), Inches(1.65), Inches(1.1), C_MID)
        add_textbox(slide, x + Inches(0.1), Inches(4.6), Inches(1.45), Inches(0.9),
                    step, font_size=10, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_textbox(slide, x + Inches(1.65), Inches(4.65), Inches(0.35), Inches(0.8),
                        ">", font_size=18, color=C_ACCENT, bold=True)
    features = [
        "目标优先级驱动", "动态预算分配", "路径代价感知", "视线遮挡判断",
        "链路质量评分", "多源融合验证", "三维态势回放", "CSV结果导出",
    ]
    for i, feat in enumerate(features):
        col = i % 4; row = i // 4
        x = Inches(1.0) + Inches(col * 3.0)
        y = Inches(5.85) + Inches(row * 0.45)
        add_textbox(slide, x, y, Inches(2.8), Inches(0.35),
                    "/ " + feat, font_size=11, color=C_WHITE)
    slide_number(slide, 6)
    return slide

# ============ SLIDE 7: SYSTEM ARCHITECTURE ============
def build_slide_07():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "系统总体架构", "System Architecture -- 轻量级前后端分离设计")
    layers = [
        ("表示层 Presentation", C_BLUE_L, [
            "参数配置界面", "策略指标卡片", "三维态势沙盘",
            "任务结果表格", "实验总结与导出",
        ]),
        ("业务控制层 Business Logic", C_MID, [
            "请求解析与路由", "实验会话管理", "步进控制与回放",
            "数据组织与JSON序列化",
        ]),
        ("仿真计算层 Simulation Engine", C_DARK, [
            "复杂环境建模", "路径规划 (栅格距离图)", "RADS/随机/全量策略调度",
            "多源感知融合与指标统计",
        ]),
    ]
    for i, (name, color, modules) in enumerate(layers):
        y = Inches(1.3) + Inches(i * 1.85)
        add_card(slide, Inches(0.5), y, Inches(3.5), Inches(1.5), color)
        add_textbox(slide, Inches(0.8), y + Inches(0.15), Inches(3.2), Inches(0.35),
                    name, font_size=15, color=C_WHITE, bold=True)
        for j, mod in enumerate(modules):
            mx = Inches(4.3) + Inches(j * 2.2)
            add_card(slide, mx, y + Inches(0.2), Inches(2.0), Inches(1.1), C_WHITE)
            add_textbox(slide, mx + Inches(0.15), y + Inches(0.35), Inches(1.7), Inches(0.7),
                        mod, font_size=10, color=C_TEXT, alignment=PP_ALIGN.CENTER)
        if i < len(layers) - 1:
            add_textbox(slide, Inches(1.8), y + Inches(1.5), Inches(1.0), Inches(0.35),
                        "V", font_size=16, color=C_ACCENT, bold=True)
    add_bg_rect(slide, Inches(0.5), Inches(6.55), Inches(12.1), Inches(0.6), C_DARK)
    add_textbox(slide, Inches(0.9), Inches(6.62), Inches(11.3), Inches(0.45),
                "技术栈: Python 3.13 + Flask + Three.js WebGL + HTML5/CSS3/JavaScript  |  HTTP + JSON  |  自研逐帧步进式会话管理",
                font_size=11, color=C_ACCENT, bold=True)
    add_figure(slide, "ch3/fig3-2_system_architecture.png", Inches(7.5), Inches(1.3), width=Inches(5.5))
    slide_number(slide, 7)
    return slide

# ============ SLIDE 8: COMPLEX SCENE MODELING ============
def build_slide_08():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "复杂场景建模", "Complex Scene Modeling -- 六要素统一仿真")
    elements = [
        ("无人机节点", "24架, 感知半径180m\n视场角120度, 速度28m/s\n初始能量100单位, 故障率6%", C_BLUE_L),
        ("动态目标", "6个, 速度14m/s\n不同优先级与不确定性\n持续运动中随机转向", RGBColor(0xE6, 0x7E, 0x22)),
        ("障碍环境", "34个建筑物, 高度36-120m\n栅格地图33x33\n影响视线和绕障路径", C_RED),
        ("通信链路", "半径300m, 丢包率4%\n延迟0-2步\n链路质量影响融合权重", RGBColor(0x8E, 0x44, 0xAD)),
        ("天气扰动", "晴空/薄雾/降雨/雷暴\n影响能见度和传感器噪声\n画面实时联动", RGBColor(0x16, 0xA0, 0x85)),
        ("地图参数", "1320x1320尺度\n栅格精度40m/格\n道路+建筑+空闲区域", C_DARK),
    ]
    for i, (name, desc, color) in enumerate(elements):
        col = i % 3; row = i // 3
        x = Inches(0.5) + Inches(col * 4.15)
        y = Inches(1.35) + Inches(row * 2.7)
        add_card(slide, x, y, Inches(3.85), Inches(2.35), C_WHITE)
        add_bg_rect(slide, x, y, Inches(3.85), Inches(0.06), color)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.35),
                    name, font_size=16, color=color, bold=True)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.4),
                    desc, font_size=11, color=C_TEXT)
    slide_number(slide, 8)
    return slide

# ============ SLIDE 9: RADS ALGORITHM CORE ============
def build_slide_09():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "RADS 动态子群选择算法", "Risk-Aware Dynamic Subgroup Selection -- 核心思想")
    add_card(slide, Inches(4.5), Inches(1.6), Inches(4.3), Inches(1.5), C_DARK)
    add_textbox(slide, Inches(4.8), Inches(1.7), Inches(3.8), Inches(0.5),
                "RADS 动态子群选择", font_size=22, color=C_ACCENT, bold=True)
    add_textbox(slide, Inches(4.8), Inches(2.25), Inches(3.8), Inches(0.6),
                "每一步重新判断:\n目标需要多少资源?\n哪些无人机更适合参与?",
                font_size=13, color=C_WHITE)
    inputs = [
        ("目标需求评估", Inches(0.5), Inches(1.5), [
            "优先级 Priority", "不确定性 Uncertainty",
            "运动强度 Motion", "目标-无人机距离",
        ]),
        ("节点能力评估", Inches(0.5), Inches(4.0), [
            "剩余能量 Energy", "链路质量 Link Quality",
            "路径代价 Path Cost", "视线条件 LOS",
            "感知覆盖能力",
        ]),
        ("环境约束感知", Inches(8.5), Inches(1.5), [
            "障碍遮挡关系", "天气扰动强度",
            "通信丢包/延迟", "故障节点排除",
        ]),
    ]
    for title, x, y, items in inputs:
        add_card(slide, x, y, Inches(3.5), Inches(2.3), C_WHITE)
        add_bg_rect(slide, x, y, Inches(3.5), Inches(0.45), C_MID)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.1), Inches(0.35),
                    title, font_size=14, color=C_WHITE, bold=True)
        for j, item in enumerate(items):
            iy = y + Inches(0.55) + Inches(j * 0.35)
            add_textbox(slide, x + Inches(0.25), iy, Inches(2.9), Inches(0.3),
                        "* " + item, font_size=11, color=C_TEXT)
    add_bg_rect(slide, Inches(4.5), Inches(3.5), Inches(4.3), Inches(1.2), RGBColor(0x14, 0x3D, 0x1A))
    add_textbox(slide, Inches(4.8), Inches(3.55), Inches(3.8), Inches(0.35),
                "输出: 无人机子群分配方案", font_size=15, color=C_GREEN, bold=True)
    add_textbox(slide, Inches(4.8), Inches(3.95), Inches(3.8), Inches(0.6),
                "确定每目标派机规模 + 成员ID\n动态预算 = f(任务压力, 系统状态)",
                font_size=11, color=C_WHITE)
    add_bg_rect(slide, Inches(0.5), Inches(6.2), Inches(12.1), Inches(0.95), C_DARK)
    add_textbox(slide, Inches(0.9), Inches(6.3), Inches(11.3), Inches(0.75),
                "对比基线: 随机派遣(随机选子群, 规模相同但成员随机) | 全量派遣(所有可用无人机全部参与)\n"
                "RADS 的核心优势: 不是派出越多越好, 而是为每个目标找到最合适的感知子群",
                font_size=12, color=C_WHITE)
    slide_number(slide, 9)
    return slide

# ============ SLIDE 10: RADS SCHEDULING FLOW ============
def build_slide_10():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "RADS 调度流程与感知融合", "Scheduling Flow & Perception Fusion Pipeline")
    steps = [
        ("Step 1\n状态更新", "更新目标和无人机\n位置、能量、故障"),
        ("Step 2\n目标需求计算", "对每个目标计算优\n先级、不确定性"),
        ("Step 3\n候选效用评估", "对每对(无人机,目标)\n计算综合效用分数"),
        ("Step 4\n动态预算分配", "按任务压力确定\n每目标派机规模"),
        ("Step 5\n子群成员选择", "按效用排序\n选择最优成员组合"),
        ("Step 6\n感知融合判定", "采集观测->链路传输\n->多机融合->成功判定"),
    ]
    for i, (title, desc) in enumerate(steps):
        x = Inches(0.3) + Inches(i * 2.15)
        add_card(slide, x, Inches(1.4), Inches(1.95), Inches(1.6), C_WHITE)
        add_textbox(slide, x + Inches(0.1), Inches(1.5), Inches(1.75), Inches(0.55),
                    title, font_size=10, color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(2.1), Inches(1.75), Inches(0.7),
                    desc, font_size=9, color=C_GRAY, alignment=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_textbox(slide, x + Inches(1.95), Inches(1.8), Inches(0.25), Inches(0.5),
                        ">", font_size=16, color=C_ACCENT, bold=True)
    add_bg_rect(slide, Inches(0.5), Inches(3.3), Inches(12.1), Inches(3.8), C_DARK)
    add_textbox(slide, Inches(0.9), Inches(3.4), Inches(11), Inches(0.35),
                "感知融合与成功判定机制", font_size=16, color=C_ACCENT, bold=True)
    fusion_steps = [
        ("观测生成", "满足距离+视线+\n视场角条件后\n概率探测目标"),
        ("链路传输", "根据丢包率和延迟\n决定观测是否送达\n影响融合权重"),
        ("稳健融合", "多机观测加权融合\n抑制异常值和\n低质量链路观测"),
        ("严格成功", "确认观测数>=阈值\n融合误差<=容忍度\n位置一致性满足"),
    ]
    for i, (title, desc) in enumerate(fusion_steps):
        x = Inches(0.9) + Inches(i * 3.0)
        add_textbox(slide, x, Inches(3.85), Inches(2.6), Inches(0.3),
                    "> " + title, font_size=13, color=C_GREEN, bold=True)
        add_textbox(slide, x, Inches(4.2), Inches(2.6), Inches(0.9),
                    desc, font_size=11, color=C_WHITE)
    add_figure(slide, "ch7/fig7-2_engine_step_flow.png",
               Inches(0.9), Inches(5.15), width=Inches(11.0), height=Inches(1.7))
    slide_number(slide, 10)
    return slide

# ============ SLIDE 11: PLATFORM IMPLEMENTATION ============
def build_slide_11():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "平台实现成果", "Platform Implementation -- 核心功能展示")
    features = [
        ("参数配置面板", "无人机/目标/障碍物数量\n感知/通信半径、丢包率\n天气类型、运行模式选择"),
        ("三维态势沙盘", "WebGL实时渲染\n无人机/目标/障碍物/路径\n旋转/缩放/俯视交互"),
        ("三策略同步对比", "主视图+两个基线窗口\n三视图大对比弹窗\n同场景公平比较"),
        ("任务级分析与导出", "目标级结果表格\n视图联动高亮\nCSV结构化导出"),
    ]
    for i, (title, desc) in enumerate(features):
        col = i % 2; row = i // 2
        x = Inches(0.5) + Inches(col * 6.3)
        y = Inches(1.35) + Inches(row * 2.8)
        add_card(slide, x, y, Inches(5.95), Inches(2.45), C_WHITE)
        add_bg_rect(slide, x, y, Inches(0.06), Inches(2.45), C_ACCENT)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.3), Inches(0.35),
                    title, font_size=16, color=C_DARK, bold=True)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.65), Inches(5.3), Inches(1.5),
                    desc, font_size=12, color=C_TEXT)
    add_bg_rect(slide, Inches(0.5), Inches(6.55), Inches(12.1), Inches(0.6), C_DARK)
    add_textbox(slide, Inches(0.9), Inches(6.62), Inches(11.3), Inches(0.45),
                "前端: HTML5 + CSS3 + JavaScript + Three.js WebGL  |  后端: Python + Flask  |  仿真引擎: 自研逐步推进会话管理",
                font_size=11, color=C_ACCENT, bold=True)
    add_figure(slide, "ch7/fig7-8_scene3d_structure.png", Inches(0.5), Inches(1.35), width=Inches(3.7))
    add_figure(slide, "ch7/fig7-1_tech_stack.png", Inches(8.5), Inches(1.35), width=Inches(4.0))
    slide_number(slide, 11)
    return slide

# ============ SLIDE 12: EXPERIMENT DESIGN ============
def build_slide_12():
    slide = add_blank_slide()
    build_slide_bg(slide)
    section_header(slide, "实验设计", "Experiment Design -- 三策略 x 四组实验 x 六维评价")
    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.35),
                "对比策略", font_size=18, color=C_DARK, bold=True)
    strategies = [
        ("RADS", "动态子群选择", "按目标需求与节点效用\n每步动态调整子群", C_GREEN),
        ("随机派遣", "Random Dispatch", "子群规模与RADS相同\n但成员随机选取", C_RED),
        ("全量派遣", "Full Dispatch", "所有可用无人机\n全部参与每个目标", C_BLUE_L),
    ]
    for i, (name, en, desc, color) in enumerate(strategies):
        x = Inches(0.5) + Inches(i * 4.1)
        add_card(slide, x, Inches(1.8), Inches(3.8), Inches(1.5), C_WHITE)
        add_bg_rect(slide, x, Inches(1.8), Inches(3.8), Inches(0.06), color)
        add_textbox(slide, x + Inches(0.25), Inches(1.95), Inches(3.3), Inches(0.35),
                    name + " (" + en + ")", font_size=15, color=color, bold=True)
        add_textbox(slide, x + Inches(0.25), Inches(2.35), Inches(3.3), Inches(0.8),
                    desc, font_size=11, color=C_GRAY)
    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11), Inches(0.35),
                "实验变量 (四组)", font_size=18, color=C_DARK, bold=True)
    groups = [
        ("默认场景", "24机6目标34障碍\n晴空天气, 标准通信", "基线参照"),
        ("任务规模", "目标数6->10\n无人机数24->32", "验证负载适应性"),
        ("天气扰动", "晴空->薄雾->降雨->雷暴\n能见度与噪声变化", "验证环境鲁棒性"),
        ("通信条件", "通信半径300->400m\n丢包率0%->10%\n延迟0->2步", "验证通信敏感性"),
    ]
    for i, (name, desc, purpose) in enumerate(groups):
        x = Inches(0.5) + Inches(i * 3.1)
        add_card(slide, x, Inches(4.05), Inches(2.85), Inches(1.6), C_WHITE)
        add_textbox(slide, x + Inches(0.2), Inches(4.15), Inches(2.45), Inches(0.3),
                    name, font_size=14, color=C_DARK, bold=True)
        add_textbox(slide, x + Inches(0.2), Inches(4.5), Inches(2.45), Inches(0.7),
                    desc, font_size=10, color=C_GRAY)
        add_textbox(slide, x + Inches(0.2), Inches(5.2), Inches(2.45), Inches(0.3),
                    "目的: " + purpose, font_size=10, color=C_ACCENT, bold=True)
    add_bg_rect(slide, Inches(0.5), Inches(5.95), Inches(12.1), Inches(0.35), C_MID)
    metrics = ["严格成功率", "定位达标率", "平均定位误差(m)", "累计能耗", "平均派机数量", "平均信息时效(步)"]
    for i, m in enumerate(metrics):
        add_textbox(slide, Inches(0.7) + Inches(i * 2.1), Inches(5.95), Inches(2.0), Inches(0.35),
                    m, font_size=10, color=C_WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(6.45), Inches(11), Inches(0.7),
                "所有策略在同一场景参数下同步运行, 保证对比公平性 | 每组实验重复多次取平均值 | 还包含障碍物数量(34->50)实验验证绕障适应性",
                font_size=10, color=C_GRAY)
    slide_number(slide, 12)
    return slide
# ============ BUILD ALL SLIDES ============
print("Running slide builders...")
builders = [
    build_slide_01, build_slide_02, build_slide_03, build_slide_04,
    build_slide_05, build_slide_06, build_slide_07, build_slide_08,
    build_slide_09, build_slide_10, build_slide_11, build_slide_12,
    build_slide_13, build_slide_14, build_slide_15, build_slide_16,
    build_slide_17, build_slide_18, build_slide_19, build_slide_20,
]

for i, builder in enumerate(builders):
    print("  Slide {}/20: {}...".format(i+1, builder.__name__))
    builder()

print("\nSaving to: " + OUTPUT_PATH)
prs.save(OUTPUT_PATH)
print("Done! File size: {:.0f} KB".format(os.path.getsize(OUTPUT_PATH) / 1024))
print("Slides: {}".format(len(prs.slides)))

# ============ SLIDE 3: RESEARCH BACKGROUND ============
